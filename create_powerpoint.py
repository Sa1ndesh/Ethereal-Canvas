#!/usr/bin/env python3
"""
Ethereal Canvas PowerPoint Presentation Generator
Creates a professional .pptx file for the AI Art & NFT Platform
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_ethereal_canvas_presentation():
    # Create presentation
    prs = Presentation()
    
    # Define colors
    purple = RGBColor(102, 126, 234)  # #667eea
    blue = RGBColor(118, 75, 162)     # #764ba2
    teal = RGBColor(78, 205, 196)     # #4ecdc4
    coral = RGBColor(255, 107, 107)   # #ff6b6b
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    
    title1.text = "🎨 Ethereal Canvas"
    title1.text_frame.paragraphs[0].font.size = Pt(54)
    title1.text_frame.paragraphs[0].font.color.rgb = purple
    
    subtitle1.text = "AI Art & NFT Platform\n\n🌟 Transform your imagination into stunning AI-generated artwork and mint them as unique NFTs on the blockchain 🌟\n\nReact 19 | TypeScript | Web3 | AI-Powered | Cross-Platform"
    
    # Slide 2: The Challenge
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    
    title2.text = "🎯 The Challenge"
    content2.text = """Current Market Problems:

🎨 Limited AI Art Access
• Complex tools require technical expertise
• Multiple expensive subscriptions needed

💰 NFT Complexity  
• Minting is complicated and expensive
• Requires blockchain knowledge

📱 Platform Fragmentation
• No unified cross-platform solution
• Scattered tools and services

🔗 Web3 Barriers
• Blockchain integration is developer-only territory
• High technical barriers for creators"""
    
    # Slide 3: Our Solution
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    content3 = slide3.placeholders[1]
    
    title3.text = "💡 Our Solution"
    content3.text = """Comprehensive Platform Features:

🤖 Multi-Provider AI
• Gemini, OpenAI, Stability AI, Hugging Face
• Intelligent failover for 99.9% uptime

⛓️ One-Click NFT Minting
• Seamless blockchain integration
• MetaMask & WalletConnect support

📱 Cross-Platform Ready
• Web, Android, Desktop applications
• Unified user experience

🎨 Professional Gallery
• Portfolio management & social sharing
• NFT tracking and analytics"""
    
    # Slide 4: Technical Architecture
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    content4 = slide4.placeholders[1]
    
    title4.text = "🏗️ Technical Architecture"
    content4.text = """Frontend Layer:
• React 19 with TypeScript
• Vite build system
• Tailwind CSS + Framer Motion
• Capacitor for mobile

Backend Layer:
• Netlify Functions
• Node.js serverless
• AI API integration
• CORS & security

Blockchain Layer:
• ethers.js 6
• WalletConnect v2
• Alchemy SDK
• IPFS storage"""
    
    # Slide 5: AI Integration
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    content5 = slide5.placeholders[1]
    
    title5.text = "🤖 AI Integration"
    content5.text = """Multi-Provider AI Services:

Intelligent Failover System:
• Hugging Face (Primary - Free Tier)
• OpenAI DALL-E (Premium Quality)
• Stability AI (High Resolution)
• Replicate (Various Models)

Enhanced Features:
• AI-optimized prompt enhancement
• Real-time generation tracking
• Multiple art styles
• High-resolution exports
• 5-15 second average generation time"""
    
    # Slide 6: Web3 & NFT Integration
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    content6 = slide6.placeholders[1]
    
    title6.text = "⛓️ Web3 & NFT Integration"
    content6.text = """Blockchain Features:

Wallet Connection:
• Universal wallet connectivity
• MetaMask, WalletConnect support
• 20+ wallet providers
• Automatic network detection

NFT Minting:
• One-click minting process
• IPFS metadata storage
• Multi-chain support (Ethereum, Polygon)
• Gas optimization
• Smart contract interactions"""
    
    # Slide 7: Mobile & Cross-Platform
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    content7 = slide7.placeholders[1]
    
    title7.text = "📱 Mobile & Cross-Platform"
    content7.text = """Android Studio Integration:
• Full Capacitor integration
• Native Android UI components
• File system access
• Push notifications
• Device integration

Platform Support:
• Web application (PWA)
• Android native app
• Desktop (Electron-ready)
• iOS (future roadmap)
• Responsive design

Development Workflow:
npm run build → npx cap sync android → npx cap open android"""
    
    # Slide 8: Performance & Metrics
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    content8 = slide8.placeholders[1]
    
    title8.text = "📊 Performance & Metrics"
    content8.text = """Core Web Vitals:
• 95+ Lighthouse Score
• <2s Load Time on 3G
• 99.9% AI Service Uptime
• <500KB Bundle Size

Optimization Features:
• Code splitting by route
• Image optimization & lazy loading
• Service worker caching
• 60fps mobile animations
• Progressive Web App capabilities"""
    
    # Slide 9: Business Model
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title9 = slide9.shapes.title
    content9 = slide9.placeholders[1]
    
    title9.text = "💰 Business Model"
    content9.text = """Revenue Streams:

🆓 Freemium Model
• Basic AI generation + limited NFT mints
• Community gallery access

💎 Premium Subscriptions
• Unlimited generations + premium models
• High-resolution exports

🏪 Marketplace Revenue
• 2.5% transaction fees
• Featured listing fees

🎯 Enterprise Solutions
• White-label licensing + custom models
• API access for businesses"""
    
    # Slide 10: Market Opportunity
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    content10 = slide10.placeholders[1]
    
    title10.text = "📈 Market Opportunity"
    content10.text = """Market Size:
• $25B AI Art Market by 2030
• $231B NFT Market Projection
• 2.8M Active NFT Users

Growth Drivers:
• 50M+ content creators seeking monetization
• 70% of digital art consumption on mobile
• 300% increase in AI tool usage
• Growing Web3 mainstream adoption"""
    
    # Slide 11: Roadmap
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    title11 = slide11.shapes.title
    content11 = slide11.placeholders[1]
    
    title11.text = "🔮 Roadmap"
    content11.text = """Phase 1 ✅ (Completed):
• AI Art Generation
• Web3 Integration
• Android App
• NFT Minting

Phase 2 🚀 (Q2 2024):
• AI Video Generation
• Music Creation
• Social Features
• Advanced Editing

Phase 3 🌍 (Q3 2024):
• iOS App
• Desktop App
• Multi-Chain Support
• IPFS Integration

Phase 4 🏪 (Q4 2024):
• Built-in Marketplace
• Creator Economy
• AI Collaborations
• Custom Models"""
    
    # Slide 12: Demo Flow
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    title12 = slide12.shapes.title
    content12 = slide12.placeholders[1]
    
    title12.text = "🎯 Demo Flow"
    content12.text = """User Journey:

1. 🎨 Generate AI Art
   Enter prompt → Select style → AI creates artwork

2. 💼 Connect Wallet
   One-click MetaMask → Network detection → Balance verification

3. 🪙 Mint as NFT
   Upload to IPFS → Smart contract → Transaction confirmation

4. 🖼️ Gallery Management
   Personal collection → NFT tracking → Social sharing

Live Demo Features:
• Real-time AI generation
• Wallet integration
• NFT minting process
• Gallery showcase"""
    
    # Slide 13: Competitive Advantage
    slide13 = prs.slides.add_slide(prs.slide_layouts[1])
    title13 = slide13.shapes.title
    content13 = slide13.placeholders[1]
    
    title13.text = "🏆 Competitive Advantage"
    content13.text = """Unique Value Propositions:

🔄 Multi-Provider Reliability
• 99.9% uptime vs single-provider solutions
• Intelligent failover system

📱 True Cross-Platform
• First native mobile + web + desktop experience
• Unified user interface

⛓️ Simplified Web3
• One-click vs complex multi-step processes
• User-friendly blockchain interaction

🎨 Creator-Centric Design
• Built by creators for creators
• Intuitive UX/UI design"""
    
    # Slide 14: Call to Action
    slide14 = prs.slides.add_slide(prs.slide_layouts[1])
    title14 = slide14.shapes.title
    content14 = slide14.placeholders[1]
    
    title14.text = "🚀 Ready to Create AI Masterpieces?"
    content14.text = """Get Started Today:

🌐 Web App
• Deploy to Netlify/Vercel instantly
• Progressive Web App ready

📱 Android App
• Open in Android Studio
• Native mobile experience

⭐ GitHub Repository
• https://github.com/Sa1ndesh/Ethereal-Canvas
• MIT Open Source License

🚀 Live Demo
• Experience the platform
• Try AI art generation

Join the future of AI-powered creative expression!"""
    
    # Slide 15: Thank You
    slide15 = prs.slides.add_slide(prs.slide_layouts[0])
    title15 = slide15.shapes.title
    subtitle15 = slide15.placeholders[1]
    
    title15.text = "Thank You! 🙏"
    title15.text_frame.paragraphs[0].font.size = Pt(54)
    title15.text_frame.paragraphs[0].font.color.rgb = teal
    
    subtitle15.text = """Questions & Discussion

🎨 Ethereal Canvas
Democratizing AI Art Creation & NFT Minting

Contact Information:
• GitHub: @Sa1ndesh
• Repository: Sa1ndesh/Ethereal-Canvas
• Platform: Cross-platform AI Art Studio

Made with ❤️ for the AI Art Community"""
    
    # Save presentation
    prs.save('Ethereal_Canvas_Presentation.pptx')
    print("✅ PowerPoint presentation created successfully!")
    print("📄 File: Ethereal_Canvas_Presentation.pptx")
    print("📊 Total slides: 15")

if __name__ == "__main__":
    try:
        create_ethereal_canvas_presentation()
    except ImportError:
        print("❌ Error: python-pptx library not found")
        print("📦 Install with: pip install python-pptx")
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
